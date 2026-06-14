#!/usr/bin/env python
"""
pdf-to-video — turn a PDF or PPTX into a narrated, animated MP4.

Pipeline (run as subcommands):
  1) extract  <input.pdf|pptx> <workdir>    -> renders slide images + slides.json (text/notes)
  2) (agent writes a polished narration.json — see SKILL.md)
  3) synth    <workdir> [--voice NAME] [--rate N]  -> per-slide WAV via Windows TTS + timing.json
  4) compose  <workdir> <output.mp4> [--no-captions]  -> Ken Burns + fades + captions, concatenated

Offline by design: visuals via PyMuPDF/Pillow, voice via Windows System.Speech,
encoding via the ffmpeg bundled with imageio-ffmpeg. No network required.
"""
import argparse
import json
import math
import os
import subprocess
import sys
import wave

HERE = os.path.dirname(os.path.abspath(__file__))
TTS_PS1 = os.path.join(HERE, "tts_win.ps1")
W, H, FPS = 1920, 1080, 30
BG = "0x0D1117"  # dark base-900


# ----------------------------- helpers -----------------------------
def ffmpeg_exe():
    import imageio_ffmpeg
    return imageio_ffmpeg.get_ffmpeg_exe()


def find_font(bold=False):
    cands = (
        ["C:/Windows/Fonts/segoeuib.ttf", "C:/Windows/Fonts/arialbd.ttf"]
        if bold
        else ["C:/Windows/Fonts/segoeui.ttf", "C:/Windows/Fonts/arial.ttf"]
    )
    for c in cands:
        if os.path.exists(c):
            return c
    return cands[0]


def wav_duration(path):
    with wave.open(path, "rb") as w:
        return w.getnframes() / float(w.getframerate())


def wrap(text, width=64):
    words, lines, cur = text.split(), [], ""
    for w in words:
        if len(cur) + len(w) + 1 <= width:
            cur = (cur + " " + w).strip()
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


# ----------------------------- extract -----------------------------
def render_pdf(input_path, slides_dir):
    import fitz
    doc = fitz.open(input_path)
    slides = []
    for i, page in enumerate(doc, start=1):
        pix = page.get_pixmap(matrix=fitz.Matrix(2.2, 2.2))
        img = os.path.join(slides_dir, f"slide_{i:03d}.png")
        pix.save(img)
        slides.append(
            {"index": i, "image": f"slides/slide_{i:03d}.png",
             "text": page.get_text("text").strip(), "notes": ""}
        )
    doc.close()
    return slides


def render_text_slide(title, body_lines, notes, out_path):
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (W, H), (13, 17, 23))
    d = ImageDraw.Draw(img)
    # subtle top accent bar
    d.rectangle([0, 0, W, 8], fill=(0, 212, 255))
    try:
        f_title = ImageFont.truetype(find_font(bold=True), 64)
        f_body = ImageFont.truetype(find_font(), 38)
    except Exception:
        f_title = ImageFont.load_default()
        f_body = ImageFont.load_default()
    margin = 120
    y = 140
    for tl in wrap(title or "", 34)[:3]:
        d.text((margin, y), tl, font=f_title, fill=(0, 212, 255))
        y += 84
    y += 30
    for line in body_lines:
        for bl in wrap(line, 70):
            if y > H - 120:
                break
            d.text((margin, y), "•  " + bl if line == bl.split("\n")[0] else bl,
                   font=f_body, fill=(230, 237, 243))
            y += 56
    img.save(out_path)


def render_pptx(input_path, slides_dir):
    from pptx import Presentation
    prs = Presentation(input_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        bullets = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt:
                continue
            if shape == slide.shapes.title or (not title and shape.has_text_frame):
                if not title:
                    title = txt
                    continue
            bullets.extend([l for l in txt.split("\n") if l.strip()])
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        img = os.path.join(slides_dir, f"slide_{i:03d}.png")
        render_text_slide(title, bullets, notes, img)
        body_text = (title + "\n" + "\n".join(bullets)).strip()
        slides.append(
            {"index": i, "image": f"slides/slide_{i:03d}.png",
             "text": body_text, "notes": notes}
        )
    return slides


def try_libreoffice_to_pdf(input_path, workdir):
    soffice = None
    for c in ["soffice", r"C:\Program Files\LibreOffice\program\soffice.exe"]:
        from shutil import which
        if which(c) or os.path.exists(c):
            soffice = c
            break
    if not soffice:
        return None
    subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                    "--outdir", workdir, input_path], check=True)
    base = os.path.splitext(os.path.basename(input_path))[0]
    pdf = os.path.join(workdir, base + ".pdf")
    return pdf if os.path.exists(pdf) else None


def cmd_extract(args):
    workdir = args.workdir
    slides_dir = os.path.join(workdir, "slides")
    os.makedirs(slides_dir, exist_ok=True)
    ext = os.path.splitext(args.input)[1].lower()

    if ext == ".pdf":
        slides = render_pdf(args.input, slides_dir)
    elif ext in (".pptx", ".ppt"):
        pdf = try_libreoffice_to_pdf(args.input, workdir)
        if pdf:
            slides = render_pdf(pdf, slides_dir)
            # also attach notes from pptx
            try:
                from pptx import Presentation
                prs = Presentation(args.input)
                for s, slide in zip(slides, prs.slides):
                    if slide.has_notes_slide:
                        s["notes"] = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass
        else:
            print("[info] LibreOffice not found — rendering slides from text via Pillow.")
            slides = render_pptx(args.input, slides_dir)
    else:
        sys.exit(f"Unsupported input: {ext} (use .pdf or .pptx)")

    out = {"input": os.path.abspath(args.input), "slides": slides}
    with open(os.path.join(workdir, "slides.json"), "w", encoding="utf-8") as f:
        json.dump(out, f, indent=2, ensure_ascii=False)
    print(f"Extracted {len(slides)} slides -> {os.path.join(workdir, 'slides.json')}")


# ----------------------------- synth -----------------------------
def load_narration(workdir):
    npath = os.path.join(workdir, "narration.json")
    if os.path.exists(npath):
        with open(npath, encoding="utf-8-sig") as f:
            return {s["index"]: s["text"] for s in json.load(f)["slides"]}
    # fallback: notes else text
    with open(os.path.join(workdir, "slides.json"), encoding="utf-8-sig") as f:
        data = json.load(f)
    out = {}
    for s in data["slides"]:
        out[s["index"]] = (s.get("notes") or s.get("text") or "").strip() or f"Slide {s['index']}."
    return out


def cmd_synth(args):
    workdir = args.workdir
    audio_dir = os.path.join(workdir, "audio")
    os.makedirs(audio_dir, exist_ok=True)
    narration = load_narration(workdir)
    timing = []
    for idx in sorted(narration):
        text = narration[idx].strip() or f"Slide {idx}."
        txt_path = os.path.join(audio_dir, f"slide_{idx:03d}.txt")
        wav_path = os.path.join(audio_dir, f"slide_{idx:03d}.wav")
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(text)
        subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-File", TTS_PS1,
             "-TextFile", txt_path, "-OutWav", wav_path, "-Voice", args.voice,
             "-Rate", str(args.rate)],
            check=True,
        )
        dur = wav_duration(wav_path)
        timing.append({"index": idx, "audio": f"audio/slide_{idx:03d}.wav",
                       "duration": round(dur, 3)})
        print(f"  slide {idx}: {dur:.1f}s  ({args.voice})")
    with open(os.path.join(workdir, "timing.json"), "w", encoding="utf-8") as f:
        json.dump({"slides": timing}, f, indent=2)
    print(f"Synthesised {len(timing)} narration clips.")


# ----------------------------- compose -----------------------------
def build_segment(ff, workdir, slide_img, wav, dur, caption_text, seg_out, captions):
    pad = 0.6
    total = dur + pad
    frames = int(math.ceil(total * FPS))
    fade_out = max(total - 0.4, 0.1)
    vf = (
        f"scale={W}:{H}:force_original_aspect_ratio=decrease,"
        f"pad={W}:{H}:(ow-iw)/2:(oh-ih)/2:color={BG},setsar=1,"
        f"zoompan=z='min(zoom+0.0006,1.08)':x='iw/2-(iw/zoom/2)':y='ih/2-(ih/zoom/2)':"
        f"d={frames}:s={W}x{H}:fps={FPS},"
        f"fade=t=in:st=0:d=0.4,fade=t=out:st={fade_out:.2f}:d=0.4"
    )
    if captions and caption_text.strip():
        cap_path = os.path.join(workdir, "audio", os.path.basename(seg_out) + ".cap.txt")
        wrapped = "\n".join(wrap(caption_text.strip(), 64)[:3])
        with open(cap_path, "w", encoding="utf-8") as f:
            f.write(wrapped)
        rel_cap = cap_path.replace("\\", "/").replace(":", "\\:")
        font = find_font().replace("\\", "/").replace(":", "\\:")
        vf += (
            f",drawtext=fontfile='{font}':textfile='{rel_cap}':reload=0:"
            f"fontcolor=white:fontsize=40:line_spacing=10:"
            f"box=1:boxcolor=0x0D1117C0:boxborderw=26:"
            f"x=(w-text_w)/2:y=h-text_h-72:enable='between(t,0.3,{total:.2f})'"
        )
    cmd = [
        ff, "-y", "-loop", "1", "-i", slide_img, "-i", wav,
        "-filter_complex", f"[0:v]{vf}[v]",
        "-map", "[v]", "-map", "1:a",
        "-t", f"{total:.2f}", "-r", str(FPS),
        "-c:v", "libx264", "-preset", "medium", "-crf", "20", "-pix_fmt", "yuv420p",
        "-c:a", "aac", "-b:a", "160k", "-ar", "44100",
        seg_out,
    ]
    subprocess.run(cmd, check=True, capture_output=True)


def cmd_compose(args):
    workdir = args.workdir
    ff = ffmpeg_exe()
    with open(os.path.join(workdir, "slides.json"), encoding="utf-8-sig") as f:
        slides = {s["index"]: s for s in json.load(f)["slides"]}
    with open(os.path.join(workdir, "timing.json"), encoding="utf-8-sig") as f:
        timing = {t["index"]: t for t in json.load(f)["slides"]}
    narration = load_narration(workdir)

    seg_dir = os.path.join(workdir, "segments")
    os.makedirs(seg_dir, exist_ok=True)
    seg_files = []
    for idx in sorted(timing):
        img = os.path.join(workdir, slides[idx]["image"])
        wav = os.path.join(workdir, timing[idx]["audio"])
        seg = os.path.join(seg_dir, f"seg_{idx:03d}.mp4")
        print(f"  rendering segment {idx}/{len(timing)} ...")
        build_segment(ff, workdir, img, wav, timing[idx]["duration"],
                      narration.get(idx, ""), seg, not args.no_captions)
        seg_files.append(seg)

    listfile = os.path.join(workdir, "concat.txt")
    with open(listfile, "w", encoding="utf-8") as f:
        for s in seg_files:
            f.write(f"file '{s.replace(os.sep, '/')}'\n")
    subprocess.run(
        [ff, "-y", "-f", "concat", "-safe", "0", "-i", listfile,
         "-c", "copy", args.output],
        check=True, capture_output=True,
    )
    print(f"\nDONE -> {args.output}")


# ----------------------------- cli -----------------------------
def main():
    p = argparse.ArgumentParser(description="PDF/PPTX -> narrated animated MP4")
    sub = p.add_subparsers(dest="cmd", required=True)

    pe = sub.add_parser("extract"); pe.add_argument("input"); pe.add_argument("workdir")
    pe.set_defaults(func=cmd_extract)

    ps = sub.add_parser("synth"); ps.add_argument("workdir")
    ps.add_argument("--voice", default="Microsoft Zira Desktop")
    ps.add_argument("--rate", type=int, default=-1)
    ps.set_defaults(func=cmd_synth)

    pc = sub.add_parser("compose"); pc.add_argument("workdir"); pc.add_argument("output")
    pc.add_argument("--no-captions", action="store_true")
    pc.set_defaults(func=cmd_compose)

    args = p.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
